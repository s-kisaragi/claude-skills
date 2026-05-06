#!/usr/bin/env python3
"""
Consulting-grade PPTX generation engine.
Accepts JSON input describing presentation structure, outputs a styled PPTX file.

Usage:
    python generate_pptx.py input.json [output.pptx] [--style mckinsey|accenture]
"""

import json
import sys
import math
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# Design Tokens
# ============================================================

@dataclass
class ColorPalette:
    primary: RGBColor
    primary_light: RGBColor
    accent1: RGBColor
    accent2: RGBColor
    accent3: RGBColor
    text_dark: RGBColor
    text_mid: RGBColor
    text_light: RGBColor
    bg_light: RGBColor
    positive: RGBColor
    negative: RGBColor
    gridline: RGBColor
    separator: RGBColor
    white: RGBColor = field(default_factory=lambda: RGBColor(0xFF, 0xFF, 0xFF))

    @property
    def chart_colors(self):
        return [self.primary, self.accent1, self.accent2, self.accent3, self.text_mid]


# Colors verified from PDF analysis (PyMuPDF extraction)
MCKINSEY_COLORS = ColorPalette(
    primary=RGBColor(0x04, 0x1C, 0x2C),       # #041C2C — full-bleed dividers
    primary_light=RGBColor(0x00, 0x58, 0x7A),  # #00587A — header bars (JP deck)
    accent1=RGBColor(0x00, 0xA9, 0xF4),        # #00A9F4 — chart accent, bright cyan
    accent2=RGBColor(0x00, 0x96, 0xD0),        # #0096D0 — secondary blue
    accent3=RGBColor(0xD4, 0xA8, 0x43),        # gold accent
    text_dark=RGBColor(0x00, 0x00, 0x00),      # #000000 — pure black body text
    text_mid=RGBColor(0x7F, 0x7F, 0x7F),       # #7F7F7F — secondary text
    text_light=RGBColor(0x89, 0x89, 0x89),     # #898989 — footer text
    bg_light=RGBColor(0xE0, 0xE0, 0xE0),       # #E0E0E0 — card/table bg
    positive=RGBColor(0x4A, 0x8C, 0x2A),
    negative=RGBColor(0xD0, 0x4A, 0x3E),
    gridline=RGBColor(0xD0, 0xD0, 0xD0),       # #D0D0D0 — chart grid
    separator=RGBColor(0xBE, 0xBF, 0xBE),      # #BEBFBE — divider lines
)

# Colors verified from 050113.pdf analysis (PyMuPDF extraction)
ACCENTURE_COLORS = ColorPalette(
    primary=RGBColor(0xA0, 0x00, 0xFF),        # #A000FF — Accenture Purple (verified)
    primary_light=RGBColor(0x75, 0x00, 0xC0),  # #7500C0 — dark purple emphasis
    accent1=RGBColor(0xB4, 0x40, 0xFF),        # #B440FF — callout/annotation purple
    accent2=RGBColor(0xEB, 0xCC, 0xFF),        # #EBCCFF — light purple highlight fill
    accent3=RGBColor(0xE6, 0xBE, 0xFF),        # #E6BEFF — text highlight background
    text_dark=RGBColor(0x00, 0x00, 0x00),      # #000000 — body text, titles
    text_mid=RGBColor(0x7F, 0x7F, 0x7F),       # #7F7F7F — source attribution
    text_light=RGBColor(0x91, 0x91, 0x91),     # #919191 — footer/copyright (verified)
    bg_light=RGBColor(0xF1, 0xF2, 0xF1),       # #F1F2F1 — card backgrounds
    positive=RGBColor(0x3E, 0x76, 0x2B),       # #3E762B — green
    negative=RGBColor(0xFF, 0x44, 0x44),
    gridline=RGBColor(0xD8, 0xD9, 0xD8),       # #D8D9D8
    separator=RGBColor(0x7E, 0x7F, 0x7E),      # #7E7F7E — badge/border gray
)


@dataclass
class Typography:
    title_font: str
    body_font: str
    jp_body_font: str          # Japanese body font
    title_size: Pt
    subtitle_size: Pt          # Also used as message_line size
    body_size: Pt
    bullet_l2_size: Pt
    chart_label_size: Pt
    exhibit_size: Pt
    source_size: Pt
    cover_title_size: Pt       # Cover/divider title size
    footer_size: Pt            # Footer firm name size
    subhead_size: Pt = None    # Subhead (section label) size
    message_size: Pt = None    # Message line size


# Font sizes verified from PDF analysis
MCKINSEY_TYPO = Typography(
    title_font="Georgia",      # Serif for titles (Bower substitute)
    body_font="Arial",         # Sans-serif for body/data
    jp_body_font="MS PGothic", # Japanese body (verified from PDF)
    title_size=Pt(25),         # 25pt — verified slide title size
    subtitle_size=Pt(18),      # 18pt — body text on content slides
    body_size=Pt(14),          # 14pt — bullet text (verified)
    bullet_l2_size=Pt(12),     # 12pt — sub-bullet
    chart_label_size=Pt(10),   # ~10pt — chart data labels
    exhibit_size=Pt(9),        # 9pt — "Exhibit N" label
    source_size=Pt(8),         # 7.9pt — source citations
    cover_title_size=Pt(48),   # 48pt — cover/divider title (verified)
    footer_size=Pt(9),         # 9.1pt — "McKinsey & Company"
)

# Font sizes verified from 050113.pdf analysis
ACCENTURE_TYPO = Typography(
    title_font="Meiryo UI",    # MeiryoUI-Bold for all titles (verified from 050113)
    body_font="Meiryo UI",     # MeiryoUI for body (verified from 050113)
    jp_body_font="Meiryo UI",  # Japanese body (MeiryoUI verified, dominant font)
    title_size=Pt(32),         # 32pt — slide title (verified from 050113)
    subtitle_size=Pt(18),      # 18pt — message line (verified from 050113)
    body_size=Pt(16),          # 16pt — body text in charts/bullets (verified)
    bullet_l2_size=Pt(14),     # 14pt — secondary body text (verified)
    chart_label_size=Pt(12),   # 12pt — chart annotations (verified)
    exhibit_size=Pt(9),        # 9pt
    source_size=Pt(8),         # 8pt — footer/copyright (Graphik-Regular italic)
    cover_title_size=Pt(32),   # 32pt — cover title (verified from 050113)
    footer_size=Pt(8),         # 8pt — copyright footer (verified from 050113)
    subhead_size=Pt(18),       # 18pt — section label subhead (verified from 050113)
    message_size=Pt(18),       # 18pt — message line (verified from 050113)
)

# Layout constants (16:9) — 960x540pt = 13.333x7.5in
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_L = Inches(0.61)       # ~43.7pt (verified from McKinsey PDF)
MARGIN_R = Inches(0.61)
HEADER_H = Inches(0.85)
CONTENT_TOP = Inches(1.15)
CONTENT_W = Inches(11.87)     # ~854pt chart content width (verified)
CONTENT_H = Inches(5.3)
FOOTER_Y = Inches(7.11)       # ~512pt from top (verified)
PAGE_NUM_X = Inches(12.5)
PAGE_NUM_Y = Inches(7.11)

# Accenture template layout indices (matching create_accenture_template.py)
ACCENTURE_LAYOUT_MAP = {
    "title": 0,               # Cover
    "content": 1,             # Content (subhead/title/message/body)
    "section_divider": 2,     # Section Divider (purple bg)
    "two_column": 3,          # Two Column
    "content_with_panel": 4,  # Content + Right Info Panel
    "agenda": 5,              # Agenda
    "chart": 6,               # Blank with Header
    "tab_content": 7,         # Tab Bar Content
    "step_content": 8,        # Numbered Step Content
    "table": 1,               # Content
    "takeaways": 1,           # Content
    "thank_you": 2,           # Section Divider (purple bg)
    "process_flow": 6,        # Blank with Header
    "matrix": 6,              # Blank with Header
    "kpi_dashboard": 6,       # Blank with Header
    "waterfall": 6,           # Blank with Header
    "appendix_divider": 2,    # Section Divider
    "three_column": 1,        # Content
}

# ============================================================
# Core Builder
# ============================================================

class PptxBuilder:
    def __init__(self, style: str = "mckinsey"):
        self.style = style
        self.colors = MCKINSEY_COLORS if style == "mckinsey" else ACCENTURE_COLORS
        self.typo = MCKINSEY_TYPO if style == "mckinsey" else ACCENTURE_TYPO
        self.page_num = 0
        self.exhibit_counter = 0

        # Load Accenture template if available
        self.uses_template = False
        if style == "accenture":
            template_path = Path(__file__).parent.parent / "assets" / "accenture_template.pptx"
            if template_path.exists():
                self.prs = Presentation(str(template_path))
                self.uses_template = True
            else:
                self.prs = Presentation()
                self.prs.slide_width = SLIDE_WIDTH
                self.prs.slide_height = SLIDE_HEIGHT
        else:
            self.prs = Presentation()
            self.prs.slide_width = SLIDE_WIDTH
            self.prs.slide_height = SLIDE_HEIGHT

    def build(self, deck: dict, output_path: str):
        """Build entire presentation from deck dict."""
        slides = deck.get("slides", [])
        for slide_data in slides:
            self.page_num += 1
            slide_type = slide_data.get("slide_type", "content")
            renderer = getattr(self, f"_render_{slide_type}", None)
            if renderer:
                slide = self._add_slide(slide_type)
                renderer(slide, slide_data)
            else:
                slide = self._add_slide("content")
                self._render_content(slide, slide_data)
        self.prs.save(output_path)
        return output_path

    def _add_slide(self, slide_type: str = "content"):
        """Add slide using appropriate layout for the style."""
        if self.uses_template:
            layout_idx = ACCENTURE_LAYOUT_MAP.get(slide_type, 1)
            layout = self.prs.slide_layouts[layout_idx]
        else:
            layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(layout)

    # ── Placeholder Helpers (template mode) ──

    def _get_placeholder(self, slide, idx):
        """Get placeholder by idx, returns None if not found."""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
        return None

    def _set_placeholder_text(self, slide, idx, text, font_name=None, font_size=None,
                              bold=None, color=None, alignment=None, italic=None):
        """Populate a placeholder with text and optional formatting."""
        ph = self._get_placeholder(slide, idx)
        if ph is None or not text:
            return None
        tf = ph.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        if font_name:
            p.font.name = font_name
        if font_size:
            p.font.size = font_size
        if bold is not None:
            p.font.bold = bold
        if color:
            p.font.color.rgb = color
        if alignment:
            p.alignment = alignment
        if italic is not None:
            p.font.italic = italic
        return ph

    def _remove_placeholder(self, slide, idx):
        """Remove an unused placeholder from the slide."""
        ph = self._get_placeholder(slide, idx)
        if ph:
            sp = ph._element
            sp.getparent().remove(sp)

    # ── Common Elements ──

    def _add_header_bar(self, slide, title: str, subtitle: str = None,
                        subhead: str = None, message_line: str = None):
        """Add styled header area. Accenture 050113 5-zone: subhead→title→message→separator."""
        if self.style == "accenture":
            msg = message_line or subtitle  # fallback to subtitle for compatibility
            if self.uses_template:
                # Template mode: populate layout placeholders
                if subhead:
                    self._set_placeholder_text(slide, 10, subhead,
                        font_name=self.typo.title_font, font_size=self.typo.subhead_size,
                        bold=True, color=self.colors.text_dark)
                else:
                    self._remove_placeholder(slide, 10)
                self._set_placeholder_text(slide, 0, title,
                    font_name=self.typo.title_font, font_size=self.typo.title_size,
                    bold=True, color=self.colors.text_dark)
                if msg:
                    self._set_placeholder_text(slide, 11, msg,
                        font_name=self.typo.title_font, font_size=self.typo.message_size,
                        bold=True, color=self.colors.text_dark)
                else:
                    self._remove_placeholder(slide, 11)
                # Separator line is on the layout — no manual creation needed
            else:
                # Fallback: programmatic shapes (no template)
                if subhead:
                    sh_txb = slide.shapes.add_textbox(
                        Pt(35), Pt(12), Inches(12.0), Pt(22)
                    )
                    sh_tf = sh_txb.text_frame
                    sh_tf.word_wrap = False
                    sh_p = sh_tf.paragraphs[0]
                    sh_p.text = subhead
                    sh_p.font.name = self.typo.title_font
                    sh_p.font.size = self.typo.subhead_size
                    sh_p.font.bold = True
                    sh_p.font.color.rgb = self.colors.text_dark
                    sh_p.alignment = PP_ALIGN.LEFT
                txb = slide.shapes.add_textbox(
                    Pt(35), Pt(32), Inches(12.0), Pt(40)
                )
                tf = txb.text_frame
                tf.word_wrap = True
                tf.margin_top = 0
                tf.margin_bottom = 0
                p = tf.paragraphs[0]
                p.text = title
                p.font.name = self.typo.title_font
                p.font.size = self.typo.title_size
                p.font.bold = True
                p.font.color.rgb = self.colors.text_dark
                p.alignment = PP_ALIGN.LEFT
                if msg:
                    msg_txb = slide.shapes.add_textbox(
                        Pt(35), Pt(73), Inches(12.0), Pt(44)
                    )
                    msg_tf = msg_txb.text_frame
                    msg_tf.word_wrap = True
                    msg_tf.margin_top = 0
                    msg_tf.margin_bottom = 0
                    msg_p = msg_tf.paragraphs[0]
                    msg_p.text = msg
                    msg_p.font.name = self.typo.title_font
                    msg_p.font.size = self.typo.message_size
                    msg_p.font.bold = True
                    msg_p.font.color.rgb = self.colors.text_dark
                    msg_p.alignment = PP_ALIGN.LEFT
                sep_y = Pt(138) if msg else Pt(75)
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Pt(35), sep_y,
                    Pt(877), Pt(0.5)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = self.colors.text_dark
                line.line.fill.background()
        else:
            # McKinsey: full-width dark blue header bar
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_H
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = self.colors.primary
            bar.line.fill.background()

            txb = slide.shapes.add_textbox(
                MARGIN_L, Inches(0.12), Inches(12.0), Inches(0.65)
            )
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.typo.title_font
            p.font.size = self.typo.title_size
            p.font.bold = True
            p.font.color.rgb = self.colors.white
            p.alignment = PP_ALIGN.LEFT
            tf.auto_size = None

    def _add_exhibit_label(self, slide, num: int = None):
        """Add exhibit label below header."""
        if num is None:
            self.exhibit_counter += 1
            num = self.exhibit_counter
        txb = slide.shapes.add_textbox(MARGIN_L, Inches(1.0), Inches(2.0), Inches(0.25))
        tf = txb.text_frame
        p = tf.paragraphs[0]
        p.text = f"Exhibit {num}"
        p.font.name = self.typo.body_font
        p.font.size = self.typo.exhibit_size
        p.font.bold = True
        p.font.color.rgb = self.colors.primary

    def _add_source(self, slide, source: str):
        """Add source citation at bottom-left."""
        if not source:
            return
        txb = slide.shapes.add_textbox(MARGIN_L, FOOTER_Y, Inches(8.0), Inches(0.25))
        tf = txb.text_frame
        p = tf.paragraphs[0]
        p.text = f"Source: {source}"
        p.font.name = self.typo.body_font
        p.font.size = self.typo.source_size
        p.font.color.rgb = self.colors.text_light

    def _add_footer_separator(self, slide):
        """Add thin separator line above footer area (McKinsey only; Accenture MCP deck has none)."""
        if self.style == "accenture":
            return  # MCP deck has no footer separator line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN_L, FOOTER_Y - Inches(0.08),
            CONTENT_W, Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.colors.separator
        line.line.fill.background()

    def _add_firm_footer(self, slide):
        """Add firm name footer. Accenture 050113: copyright left + page number right at Y=519."""
        if self.style == "accenture":
            if self.uses_template:
                return  # Copyright inherited from slide master
            # Fallback: programmatic copyright
            copy_txb = slide.shapes.add_textbox(
                Pt(35), Pt(519), Inches(6.0), Pt(12)
            )
            copy_tf = copy_txb.text_frame
            copy_tf.margin_top = 0
            copy_tf.margin_bottom = 0
            copy_p = copy_tf.paragraphs[0]
            copy_p.text = "Copyright \u00a9 2026 Accenture. All rights reserved."
            copy_p.font.name = "Graphik"
            copy_p.font.size = self.typo.footer_size
            copy_p.font.color.rgb = self.colors.text_light
            copy_p.font.italic = True
            copy_p.alignment = PP_ALIGN.LEFT
        else:
            text = "McKinsey & Company"
            txb = slide.shapes.add_textbox(
                Inches(9.0), PAGE_NUM_Y, Inches(3.5), Inches(0.25)
            )
            tf = txb.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.name = self.typo.body_font
            p.font.size = self.typo.footer_size
            p.font.color.rgb = self.colors.text_light
            p.alignment = PP_ALIGN.RIGHT

    def _add_page_number(self, slide):
        """Add page number. Accenture 050113: right side Y=519. McKinsey: bottom-right."""
        if self.style == "accenture":
            # Always use programmatic textbox (sldNum placeholders don't auto-inherit)
            if self.uses_template:
                self._remove_placeholder(slide, 12)  # remove template placeholder if present
            txb = slide.shapes.add_textbox(Pt(890), Pt(519), Pt(40), Pt(12))
            tf = txb.text_frame
            tf.margin_top = 0
            tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = str(self.page_num)
            p.font.name = "Graphik"
            p.font.size = Pt(8)
            p.font.color.rgb = self.colors.text_light
            p.font.italic = True
            p.alignment = PP_ALIGN.RIGHT
        else:
            txb = slide.shapes.add_textbox(PAGE_NUM_X, PAGE_NUM_Y, Inches(0.5), Inches(0.25))
            tf = txb.text_frame
            p = tf.paragraphs[0]
            p.text = str(self.page_num)
            p.font.name = self.typo.body_font
            p.font.size = self.typo.source_size
            p.font.color.rgb = self.colors.text_light
            p.alignment = PP_ALIGN.RIGHT

    def _set_cell_border(self, cell, bottom_color=None, bottom_width=None):
        """Set table cell border via XML (python-pptx lacks direct API)."""
        tc = cell._tc
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn('a:tcPr'))
        if bottom_color:
            lnB = etree.SubElement(tcPr, qn('a:lnB'))
            lnB.set('w', str(int(bottom_width) if bottom_width else '12700'))
            solidFill = etree.SubElement(lnB, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            srgbClr.set('val', f"{bottom_color[0]:02X}{bottom_color[1]:02X}{bottom_color[2]:02X}")

    def _add_accenture_gt_mark(self, slide):
        """Add Accenture > brand mark — skipped for 050113-style (no > in footer)."""
        if self.style != "accenture":
            return
        # 050113 style has no ">" chevron in the footer — intentionally blank
        pass

    def _add_bullets(self, slide, bullets, left, top, width, height):
        """Add bullet list to slide."""
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf = txb.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(bullet, dict):
                p.text = f"  {bullet['text']}"
                p.font.name = self.typo.body_font
                p.font.size = self.typo.body_size
                p.font.color.rgb = self.colors.text_dark
                p.space_after = Pt(4)
                for sub in bullet.get("sub", []):
                    sp = tf.add_paragraph()
                    sp.text = f"    – {sub}"
                    sp.font.name = self.typo.body_font
                    sp.font.size = self.typo.bullet_l2_size
                    sp.font.color.rgb = self.colors.text_mid
                    sp.space_after = Pt(2)
                    sp.level = 1
            else:
                p.text = f"  {bullet}"
                p.font.name = self.typo.body_font
                p.font.size = self.typo.body_size
                p.font.color.rgb = self.colors.text_dark
                p.space_after = Pt(4)
        return txb

    def _style_common(self, slide, data: dict):
        """Apply common elements (source, separator, firm name, page number, exhibit, brand)."""
        self._add_footer_separator(slide)
        if data.get("source"):
            self._add_source(slide, data["source"])
        self._add_firm_footer(slide)
        self._add_page_number(slide)
        if data.get("exhibit_number") is not None:
            self._add_exhibit_label(slide, data["exhibit_number"])
        self._add_accenture_gt_mark(slide)

    # ── Slide Renderers ──

    def _render_title(self, slide, data: dict):
        """Title / cover slide."""
        if self.style == "accenture":
            self._render_title_accenture(slide, data)
        else:
            self._render_title_mckinsey(slide, data)

    def _render_title_mckinsey(self, slide, data: dict):
        """McKinsey title: solid dark blue, serif title, conservative."""
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors.primary

        txb = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.0), Inches(10.0), Inches(2.0)
        )
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.name = self.typo.title_font
        p.font.size = self.typo.cover_title_size  # 48pt (verified)
        p.font.bold = True
        p.font.color.rgb = self.colors.white
        p.alignment = PP_ALIGN.LEFT

        if data.get("subtitle"):
            p2 = tf.add_paragraph()
            p2.text = data["subtitle"]
            p2.font.name = self.typo.body_font
            p2.font.size = Pt(18)
            p2.font.color.rgb = self.colors.white
            p2.space_before = Pt(12)

        meta_parts = []
        if data.get("date"):
            meta_parts.append(data["date"])
        if data.get("author"):
            meta_parts.append(data["author"])
        if meta_parts:
            txb2 = slide.shapes.add_textbox(
                Inches(1.0), Inches(5.5), Inches(10.0), Inches(0.5)
            )
            tf2 = txb2.text_frame
            p3 = tf2.paragraphs[0]
            p3.text = " | ".join(meta_parts)
            p3.font.name = self.typo.body_font
            p3.font.size = Pt(12)
            p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

        if data.get("confidential"):
            txb3 = slide.shapes.add_textbox(
                Inches(1.0), Inches(6.5), Inches(4.0), Inches(0.3)
            )
            tf3 = txb3.text_frame
            p4 = tf3.paragraphs[0]
            p4.text = "CONFIDENTIAL"
            p4.font.name = self.typo.body_font
            p4.font.size = Pt(8)
            p4.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            p4.font.bold = True

    def _render_title_accenture(self, slide, data: dict):
        """Accenture 050113-style title: white bg, addressee, project title, date, copyright."""
        if self.uses_template:
            # Cover layout: idx 14=addressee, 15=subtitle, 0=title, 16=doc_type, 17=date
            font = self.typo.title_font
            clr = self.colors.text_dark
            if data.get("addressee"):
                self._set_placeholder_text(slide, 14, data["addressee"],
                    font_name=font, font_size=Pt(16), color=clr)
            else:
                self._remove_placeholder(slide, 14)
            if data.get("subtitle"):
                self._set_placeholder_text(slide, 15, data["subtitle"],
                    font_name=font, font_size=Pt(18), bold=True, color=clr)
            else:
                self._remove_placeholder(slide, 15)
            self._set_placeholder_text(slide, 0, data.get("title", ""),
                font_name=font, font_size=self.typo.cover_title_size, bold=True, color=clr)
            if data.get("document_type"):
                self._set_placeholder_text(slide, 16, data["document_type"],
                    font_name=font, font_size=Pt(20), bold=True, color=clr)
            else:
                self._remove_placeholder(slide, 16)
            if data.get("date"):
                self._set_placeholder_text(slide, 17, data["date"],
                    font_name=font, font_size=Pt(20), bold=True, color=clr)
            else:
                self._remove_placeholder(slide, 17)
            # Copyright from master — no manual creation needed
            return

        # Fallback: programmatic shapes (no template)
        if data.get("addressee"):
            addr_txb = slide.shapes.add_textbox(
                Pt(35), Pt(119), Inches(8.0), Pt(22)
            )
            addr_tf = addr_txb.text_frame
            addr_p = addr_tf.paragraphs[0]
            addr_p.text = data["addressee"]
            addr_p.font.name = self.typo.title_font
            addr_p.font.size = Pt(16)
            addr_p.font.color.rgb = self.colors.text_dark

        if data.get("subtitle"):
            sub_txb = slide.shapes.add_textbox(
                Pt(35), Pt(157), Inches(11.0), Pt(24)
            )
            sub_tf = sub_txb.text_frame
            sub_p = sub_tf.paragraphs[0]
            sub_p.text = data["subtitle"]
            sub_p.font.name = self.typo.title_font
            sub_p.font.size = Pt(18)
            sub_p.font.bold = True
            sub_p.font.color.rgb = self.colors.text_dark

        txb = slide.shapes.add_textbox(
            Pt(35), Pt(190), Inches(11.0), Pt(70)
        )
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.name = self.typo.title_font
        p.font.size = self.typo.cover_title_size
        p.font.bold = True
        p.font.color.rgb = self.colors.text_dark
        p.alignment = PP_ALIGN.LEFT

        if data.get("document_type"):
            doc_txb = slide.shapes.add_textbox(
                Pt(35), Pt(284), Inches(5.0), Pt(28)
            )
            doc_tf = doc_txb.text_frame
            doc_p = doc_tf.paragraphs[0]
            doc_p.text = data["document_type"]
            doc_p.font.name = self.typo.title_font
            doc_p.font.size = Pt(20)
            doc_p.font.bold = True
            doc_p.font.color.rgb = self.colors.text_dark

        if data.get("date"):
            date_txb = slide.shapes.add_textbox(
                Pt(80), Pt(373), Inches(4.0), Pt(28)
            )
            date_tf = date_txb.text_frame
            date_p = date_tf.paragraphs[0]
            date_p.text = data["date"]
            date_p.font.name = self.typo.title_font
            date_p.font.size = Pt(20)
            date_p.font.bold = True
            date_p.font.color.rgb = self.colors.text_dark

        copy_txb = slide.shapes.add_textbox(
            Pt(35), Pt(527), Inches(6.0), Pt(12)
        )
        copy_tf = copy_txb.text_frame
        copy_tf.margin_top = 0
        copy_p = copy_tf.paragraphs[0]
        copy_p.text = "Copyright \u00a9 2026 Accenture. All rights reserved."
        copy_p.font.name = "Graphik"
        copy_p.font.size = Pt(8)
        copy_p.font.color.rgb = self.colors.text_mid
        copy_p.font.italic = True

    def _render_agenda(self, slide, data: dict):
        """Agenda slide. Accenture 050113: full-width #A000FF purple bands. McKinsey: circles."""
        items = data.get("items", [])
        active = data.get("active_index")

        if self.style == "accenture":
            if self.uses_template:
                # Use agenda layout title placeholder
                self._set_placeholder_text(slide, 0, data.get("title", "Agenda"),
                    font_name=self.typo.title_font, font_size=Pt(28),
                    bold=True, color=self.colors.text_dark)
                # Remove body placeholder — items are drawn programmatically
                self._remove_placeholder(slide, 1)
            else:
                self._add_header_bar(slide, data.get("title", "Agenda"))

            # 050113 style: full-width #A000FF purple bands
            item_start_y = Pt(165)
            item_h = Pt(32)
            item_spacing = Pt(49)

            for i, item in enumerate(items):
                y = item_start_y + i * item_spacing
                is_active = (active is not None and i == active)

                if is_active:
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Pt(0), y, Pt(960), item_h
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = self.colors.primary
                    bar.line.fill.background()

                txb = slide.shapes.add_textbox(
                    Pt(48), y + Pt(2), Pt(860), Pt(28)
                )
                tf = txb.text_frame
                tf.margin_top = 0
                tf.margin_bottom = 0
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = f"{i + 1}．{item}"
                p.font.name = self.typo.title_font
                p.font.size = Pt(24)
                p.font.bold = is_active
                p.font.color.rgb = self.colors.white if is_active else self.colors.text_dark
                p.alignment = PP_ALIGN.LEFT

            self._add_firm_footer(slide)
            self._add_page_number(slide)
        else:
            # McKinsey: numbered circles style
            self._add_header_bar(slide, data.get("title", "Agenda"))
            y = CONTENT_TOP + Inches(0.3)
            for i, item in enumerate(items):
                num_shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(1.0), y, Inches(0.5), Inches(0.5)
                )
                is_active = (active is not None and i == active)
                num_shape.fill.solid()
                num_shape.fill.fore_color.rgb = self.colors.primary if is_active else self.colors.bg_light
                num_shape.line.fill.background()
                ntf = num_shape.text_frame
                ntf.paragraphs[0].text = str(i + 1)
                ntf.paragraphs[0].font.size = Pt(14)
                ntf.paragraphs[0].font.bold = True
                ntf.paragraphs[0].font.color.rgb = self.colors.white if is_active else self.colors.text_mid
                ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
                ntf.vertical_anchor = MSO_ANCHOR.MIDDLE

                txb = slide.shapes.add_textbox(Inches(1.8), y, Inches(9.0), Inches(0.5))
                tf = txb.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = item
                p.font.name = self.typo.body_font
                p.font.size = Pt(14)
                p.font.bold = is_active
                p.font.color.rgb = self.colors.text_dark if is_active else self.colors.text_mid

                y += Inches(0.8)

            self._add_page_number(slide)

    def _render_section_divider(self, slide, data: dict):
        """Section divider with style-specific background."""
        if self.style == "accenture":
            self._render_section_divider_accenture(slide, data)
        else:
            self._render_section_divider_mckinsey(slide, data)

    def _render_section_divider_mckinsey(self, slide, data: dict):
        """McKinsey section divider: solid dark blue background."""
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.colors.primary

        if data.get("section_number") is not None:
            txb = slide.shapes.add_textbox(
                Inches(1.0), Inches(2.0), Inches(2.0), Inches(1.5)
            )
            tf = txb.text_frame
            p = tf.paragraphs[0]
            p.text = f"{data['section_number']:02d}"
            p.font.name = self.typo.title_font
            p.font.size = Pt(60)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        title_left = Inches(1.0) if not data.get("section_number") else Inches(3.5)
        txb2 = slide.shapes.add_textbox(
            title_left, Inches(2.5), Inches(9.0), Inches(2.0)
        )
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = data.get("section_title", "")
        p2.font.name = self.typo.title_font
        p2.font.size = self.typo.cover_title_size  # 48pt (verified)
        p2.font.bold = True
        p2.font.color.rgb = self.colors.white

        if data.get("subtitle"):
            p3 = tf2.add_paragraph()
            p3.text = data["subtitle"]
            p3.font.name = self.typo.body_font
            p3.font.size = Pt(16)
            p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            p3.space_before = Pt(12)

    def _render_section_divider_accenture(self, slide, data: dict):
        """Accenture 050113-style section divider: purple background, white text."""
        if self.uses_template:
            # Section divider layout: purple bg from layout, idx 10=number, 0=title, 11=subtitle
            if data.get("section_number") is not None:
                self._set_placeholder_text(slide, 10, f"{data['section_number']}.",
                    font_name=self.typo.title_font, font_size=Pt(40),
                    bold=True, color=self.colors.white)
            else:
                self._remove_placeholder(slide, 10)
            self._set_placeholder_text(slide, 0, data.get("section_title", ""),
                font_name=self.typo.title_font, font_size=Pt(32),
                bold=True, color=self.colors.white)
            if data.get("subtitle"):
                self._set_placeholder_text(slide, 11, data["subtitle"],
                    font_name=self.typo.body_font, font_size=Pt(18),
                    color=RGBColor(0xDD, 0xDD, 0xDD))
            else:
                self._remove_placeholder(slide, 11)
            return

        # Fallback: programmatic shapes (no template)
        bg_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
        )
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = self.colors.primary
        bg_rect.line.fill.background()

        if data.get("section_number") is not None:
            num_txb = slide.shapes.add_textbox(
                Pt(48), Pt(180), Pt(60), Pt(50)
            )
            num_tf = num_txb.text_frame
            num_p = num_tf.paragraphs[0]
            num_p.text = f"{data['section_number']}."
            num_p.font.name = self.typo.title_font
            num_p.font.size = Pt(40)
            num_p.font.bold = True
            num_p.font.color.rgb = self.colors.white

        txb2 = slide.shapes.add_textbox(
            Pt(48), Pt(240), Pt(860), Pt(80)
        )
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = data.get("section_title", "")
        p2.font.name = self.typo.title_font
        p2.font.size = Pt(32)
        p2.font.bold = True
        p2.font.color.rgb = self.colors.white

        if data.get("subtitle"):
            p3 = tf2.add_paragraph()
            p3.text = data["subtitle"]
            p3.font.name = self.typo.body_font
            p3.font.size = Pt(18)
            p3.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            p3.space_before = Pt(12)

    def _render_content(self, slide, data: dict):
        """Content slide with action title and bullets."""
        msg = data.get("message_line") or data.get("subtitle")
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=msg)
        # Remove body placeholder if template — we add bullets programmatically
        if self.uses_template:
            self._remove_placeholder(slide, 1)
        # Body starts after separator line — Y~155pt for Accenture with message, ~80pt without
        if self.style == "accenture":
            if msg:
                body_top = Pt(155)
            else:
                body_top = Pt(80)
        else:
            body_top = CONTENT_TOP + Inches(0.2)
        bullets = data.get("bullets", [])
        if bullets:
            body_h = Pt(355) if self.style == "accenture" else CONTENT_H - Inches(0.5)
            self._add_bullets(
                slide, bullets,
                Pt(35) if self.style == "accenture" else MARGIN_L,
                body_top, CONTENT_W, body_h
            )
        self._style_common(slide, data)

    def _render_chart(self, slide, data: dict):
        """Chart slide with single chart."""
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=data.get("message_line") or data.get("subtitle"))

        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
            "stacked_bar_100": XL_CHART_TYPE.COLUMN_STACKED_100,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "area": XL_CHART_TYPE.AREA,
        }

        ct = data.get("chart_type", "column")
        xl_type = chart_type_map.get(ct, XL_CHART_TYPE.COLUMN_CLUSTERED)

        categories = data.get("categories", [])
        series = data.get("series", {})

        chart_left = Inches(0.8)
        chart_top = CONTENT_TOP + Inches(0.3)
        chart_width = Inches(11.5)
        chart_height = Inches(4.5)

        if ct == "scatter":
            chart_data = XyChartData()
            for name, values in series.items():
                s = chart_data.add_series(name)
                for point in values:
                    s.add_data_point(point[0], point[1])
        else:
            chart_data = CategoryChartData()
            chart_data.categories = categories
            for name, values in series.items():
                chart_data.add_series(name, values)

        graphic_frame = slide.shapes.add_chart(
            xl_type, chart_left, chart_top, chart_width, chart_height, chart_data
        )
        chart = graphic_frame.chart

        # Style chart
        colors = self.colors.chart_colors
        highlight_idx = data.get("highlight_index")

        for i, s in enumerate(chart.series):
            s.format.fill.solid()
            if highlight_idx is not None and ct in ("column", "bar"):
                s.format.fill.fore_color.rgb = self.colors.gridline
                try:
                    pt = s.points[highlight_idx]
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = self.colors.primary
                except (IndexError, AttributeError):
                    pass
            else:
                s.format.fill.fore_color.rgb = colors[i % len(colors)]

        # Data labels
        nf = data.get("number_format", "#,##0")
        for plot in chart.plots:
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.font.size = self.typo.chart_label_size
            dl.font.name = self.typo.body_font
            dl.font.color.rgb = self.colors.text_dark
            dl.number_format = nf
            dl.show_value = True
            dl.show_category_name = False
            dl.show_series_name = False
            try:
                dl.label_position = XL_LABEL_POSITION.OUTSIDE_END
            except ValueError:
                pass

        # Legend
        if len(series) > 1:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)
            chart.legend.font.name = self.typo.body_font
        else:
            chart.has_legend = False

        # Axes
        try:
            va = chart.value_axis
            va.has_major_gridlines = True
            va.major_gridlines.format.line.color.rgb = self.colors.gridline
            va.major_gridlines.format.line.width = Pt(0.5)
            va.has_minor_gridlines = False
            va.tick_labels.font.size = Pt(9)
            va.tick_labels.font.color.rgb = self.colors.text_mid

            if data.get("y_axis_label"):
                va.has_title = True
                va.axis_title.text_frame.paragraphs[0].text = data["y_axis_label"]
                va.axis_title.text_frame.paragraphs[0].font.size = Pt(9)

            ca = chart.category_axis
            ca.tick_labels.font.size = Pt(9)
            ca.tick_labels.font.color.rgb = self.colors.text_mid
        except (ValueError, AttributeError):
            pass

        # Annotation callout
        if data.get("annotation"):
            ann_txb = slide.shapes.add_textbox(
                Inches(8.5), CONTENT_TOP + Inches(0.3),
                Inches(4.0), Inches(0.6)
            )
            ann_shape = ann_txb
            ann_tf = ann_txb.text_frame
            ann_tf.word_wrap = True
            p = ann_tf.paragraphs[0]
            p.text = data["annotation"]
            p.font.name = self.typo.body_font
            p.font.size = Pt(10)
            p.font.color.rgb = self.colors.primary
            # Add border via XML
            spPr = ann_txb._element.find(qn('p:spPr'))
            if spPr is None:
                spPr = etree.SubElement(ann_txb._element, qn('p:spPr'))
            ln = etree.SubElement(spPr, qn('a:ln'))
            ln.set('w', '12700')
            solidFill = etree.SubElement(ln, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            hex_color = f"{self.colors.primary[0]:02X}{self.colors.primary[1]:02X}{self.colors.primary[2]:02X}"
            srgbClr.set('val', hex_color)

        self._style_common(slide, data)

    def _render_two_column(self, slide, data: dict):
        """Two-column layout."""
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=data.get("message_line") or data.get("subtitle"))
        # Remove body placeholders — columns are drawn programmatically
        if self.uses_template:
            self._remove_placeholder(slide, 1)
            self._remove_placeholder(slide, 13)
        split = data.get("split", "equal")

        if split == "wide_left":
            lw, rw = Inches(7.2), Inches(4.6)
            rx = Inches(8.2)
        elif split == "wide_right":
            lw, rw = Inches(4.6), Inches(7.2)
            rx = Inches(5.6)
        else:
            lw, rw = Inches(5.9), Inches(5.9)
            rx = Inches(6.9)

        y = CONTENT_TOP + Inches(0.2)
        h = CONTENT_H - Inches(0.5)

        for col_data, x, w in [(data.get("left", {}), MARGIN_L, lw),
                                (data.get("right", {}), rx, rw)]:
            if col_data.get("heading"):
                txb = slide.shapes.add_textbox(x, y, w, Inches(0.4))
                tf = txb.text_frame
                p = tf.paragraphs[0]
                p.text = col_data["heading"]
                p.font.name = self.typo.body_font
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = self.colors.primary

            if col_data.get("bullets"):
                self._add_bullets(slide, col_data["bullets"], x, y + Inches(0.5), w, h - Inches(0.5))

        self._style_common(slide, data)

    def _render_three_column(self, slide, data: dict):
        """Three-column layout."""
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=data.get("message_line") or data.get("subtitle"))
        if self.uses_template:
            self._remove_placeholder(slide, 1)
        columns = data.get("columns", [])
        col_w = Inches(3.78)
        col_positions = [MARGIN_L, Inches(4.68), Inches(8.86)]
        y = CONTENT_TOP + Inches(0.2)
        h = CONTENT_H - Inches(0.5)

        for i, col in enumerate(columns[:3]):
            x = col_positions[i]

            if col.get("heading"):
                txb = slide.shapes.add_textbox(x, y, col_w, Inches(0.4))
                tf = txb.text_frame
                p = tf.paragraphs[0]
                p.text = col["heading"]
                p.font.name = self.typo.body_font
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = self.colors.primary

            if col.get("bullets"):
                self._add_bullets(slide, col["bullets"], x, y + Inches(0.5), col_w, h - Inches(0.5))

        self._style_common(slide, data)

    def _render_table(self, slide, data: dict):
        """Table slide."""
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=data.get("message_line") or data.get("subtitle"))
        if self.uses_template:
            self._remove_placeholder(slide, 1)
        headers = data.get("headers", [])
        rows = data.get("rows", [])
        n_rows = len(rows) + 1
        n_cols = len(headers)

        table_top = CONTENT_TOP + Inches(0.3)
        table_h = min(Inches(0.4) * n_rows, CONTENT_H - Inches(0.8))
        table_shape = slide.shapes.add_table(
            n_rows, n_cols, MARGIN_L, table_top, CONTENT_W, table_h
        )
        table = table_shape.table

        # Column widths
        col_widths = data.get("col_widths")
        if col_widths:
            for i, cw in enumerate(col_widths):
                table.columns[i].width = int(CONTENT_W * cw)
        else:
            w = int(CONTENT_W / n_cols)
            for i in range(n_cols):
                table.columns[i].width = w

        # Header row
        for i, h_text in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h_text
            p = cell.text_frame.paragraphs[0]
            p.font.name = self.typo.body_font
            p.font.size = Pt(10)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if self.style == "accenture":
                # Accenture 050113: #404040 dark gray header, white text
                p.font.color.rgb = self.colors.white
                p.font.size = Pt(12)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x40)  # #404040
            else:
                # McKinsey: solid primary header
                p.font.color.rgb = self.colors.white
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.colors.primary

        # Data rows
        highlight_rows = data.get("highlight_rows", [])
        acn_highlight = RGBColor(0xF3, 0xE5, 0xFF)  # light purple for Accenture
        mck_highlight = RGBColor(0xE8, 0xF0, 0xFE)  # light blue for McKinsey
        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                p.font.name = self.typo.body_font
                p.font.size = Pt(10)
                p.font.color.rgb = self.colors.text_dark
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                if r_idx in highlight_rows:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = acn_highlight if self.style == "accenture" else mck_highlight
                elif r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.colors.bg_light
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.colors.white

        self._style_common(slide, data)

    def _render_process_flow(self, slide, data: dict):
        """Process flow with horizontal boxes and arrows."""
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=data.get("message_line") or data.get("subtitle"))
        steps = data.get("steps", [])
        n = len(steps)
        if n == 0:
            return

        total_w = CONTENT_W
        box_w = Inches(2.2) if n <= 4 else Inches(1.8)
        arrow_w = Inches(0.4)
        total_needed = box_w * n + arrow_w * (n - 1)
        start_x = MARGIN_L + (total_w - total_needed) / 2
        y_center = CONTENT_TOP + CONTENT_H / 2 - Inches(0.5)
        box_h = Inches(1.6)

        for i, step in enumerate(steps):
            x = start_x + i * (box_w + arrow_w)

            # Box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y_center), int(box_w), int(box_h)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors.white
            shape.line.color.rgb = self.colors.primary
            shape.line.width = Pt(1.5)

            # Label
            tf = shape.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = step.get("label", "")
            p.font.name = self.typo.body_font
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.colors.primary
            p.alignment = PP_ALIGN.CENTER

            if step.get("description"):
                p2 = tf.add_paragraph()
                p2.text = step["description"]
                p2.font.name = self.typo.body_font
                p2.font.size = Pt(9)
                p2.font.color.rgb = self.colors.text_mid
                p2.alignment = PP_ALIGN.CENTER

            # Arrow (except after last box)
            if i < n - 1:
                arrow_x = int(x + box_w)
                arrow_y = int(y_center + box_h / 2)
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR_TYPE.STRAIGHT,
                    arrow_x, arrow_y,
                    int(arrow_x + arrow_w), arrow_y
                )
                connector.line.color.rgb = self.colors.separator
                connector.line.width = Pt(1.5)

        self._style_common(slide, data)

    def _render_matrix(self, slide, data: dict):
        """Comparison matrix with Harvey balls or RAG status."""
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=data.get("message_line") or data.get("subtitle"))
        row_headers = data.get("row_headers", [])
        col_headers = data.get("col_headers", [])
        cells = data.get("cells", [])
        cell_type = data.get("cell_type", "text")

        harvey_map = {0: "○", 1: "◔", 2: "◑", 3: "◕", 4: "●"}
        jp_eval_map = {0: "×", 1: "△", 2: "○", 3: "◎"}

        n_rows = len(row_headers) + 1
        n_cols = len(col_headers) + 1

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, MARGIN_L, CONTENT_TOP + Inches(0.3),
            CONTENT_W, min(Inches(0.45) * n_rows, CONTENT_H - Inches(0.8))
        )
        table = table_shape.table

        # Header row
        table.cell(0, 0).text = ""
        table.cell(0, 0).fill.solid()
        table.cell(0, 0).fill.fore_color.rgb = self.colors.primary
        for j, ch in enumerate(col_headers):
            cell = table.cell(0, j + 1)
            cell.text = ch
            p = cell.text_frame.paragraphs[0]
            p.font.name = self.typo.body_font
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.colors.white
            p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors.primary
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data rows
        for i, rh in enumerate(row_headers):
            cell = table.cell(i + 1, 0)
            cell.text = rh
            p = cell.text_frame.paragraphs[0]
            p.font.name = self.typo.body_font
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.colors.text_dark
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors.bg_light if i % 2 == 0 else self.colors.white

            for j, val in enumerate(cells[i] if i < len(cells) else []):
                cell = table.cell(i + 1, j + 1)
                if cell_type == "harvey_ball":
                    display = harvey_map.get(val, str(val))
                elif cell_type == "jp_eval":
                    display = jp_eval_map.get(val, str(val))
                else:
                    display = str(val)
                cell.text = display
                p = cell.text_frame.paragraphs[0]
                p.font.name = self.typo.body_font
                p.font.size = Pt(14) if cell_type in ("harvey_ball", "jp_eval") else Pt(10)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = self.colors.text_dark
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.colors.bg_light if i % 2 == 0 else self.colors.white
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        self._style_common(slide, data)

    def _render_kpi_dashboard(self, slide, data: dict):
        """KPI dashboard with metric cards."""
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=data.get("message_line") or data.get("subtitle"))
        kpis = data.get("kpis", [])
        n = len(kpis)
        if n == 0:
            return

        card_w = Inches(2.8)
        card_h = Inches(2.0)
        gap = Inches(0.3)
        total_w = card_w * n + gap * (n - 1)
        start_x = MARGIN_L + (CONTENT_W - total_w) / 2
        y = CONTENT_TOP + Inches(1.0)

        for i, kpi in enumerate(kpis):
            x = start_x + i * (card_w + gap)

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(card_w), int(card_h)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.colors.white
            card.line.color.rgb = self.colors.gridline
            card.line.width = Pt(0.75)

            # Value
            txb = slide.shapes.add_textbox(
                int(x + Inches(0.2)), int(y + Inches(0.2)),
                int(card_w - Inches(0.4)), int(Inches(0.8))
            )
            tf = txb.text_frame
            p = tf.paragraphs[0]
            p.text = str(kpi.get("value", ""))
            p.font.name = self.typo.title_font
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.colors.primary
            p.alignment = PP_ALIGN.CENTER

            # Label
            txb2 = slide.shapes.add_textbox(
                int(x + Inches(0.2)), int(y + Inches(1.0)),
                int(card_w - Inches(0.4)), int(Inches(0.4))
            )
            tf2 = txb2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = kpi.get("label", "")
            p2.font.name = self.typo.body_font
            p2.font.size = Pt(10)
            p2.font.color.rgb = self.colors.text_mid
            p2.alignment = PP_ALIGN.CENTER

            # Delta
            if kpi.get("delta"):
                txb3 = slide.shapes.add_textbox(
                    int(x + Inches(0.2)), int(y + Inches(1.4)),
                    int(card_w - Inches(0.4)), int(Inches(0.3))
                )
                tf3 = txb3.text_frame
                p3 = tf3.paragraphs[0]
                trend = kpi.get("trend", "up")
                arrow = "▲" if trend == "up" else "▼" if trend == "down" else "—"
                p3.text = f"{arrow} {kpi['delta']}"
                p3.font.name = self.typo.body_font
                p3.font.size = Pt(11)
                p3.font.bold = True
                p3.font.color.rgb = self.colors.positive if trend == "up" else self.colors.negative if trend == "down" else self.colors.text_mid
                p3.alignment = PP_ALIGN.CENTER

        self._style_common(slide, data)

    def _render_waterfall(self, slide, data: dict):
        """Waterfall chart built with rectangles."""
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=data.get("message_line") or data.get("subtitle"))
        items = data.get("items", [])
        n = len(items)
        if n == 0:
            return

        nf = data.get("number_format", "#,##0")
        chart_left = Inches(1.0)
        chart_top = CONTENT_TOP + Inches(0.5)
        chart_w = Inches(11.0)
        chart_h = Inches(4.5)

        values = [item["value"] for item in items]
        max_cumulative = 0
        min_cumulative = 0
        cumulative = 0
        for item in items:
            if item.get("type") == "total":
                cumulative = item["value"]
            else:
                cumulative += item["value"]
            max_cumulative = max(max_cumulative, cumulative)
            min_cumulative = min(min_cumulative, cumulative)

        value_range = max_cumulative - min_cumulative
        if value_range == 0:
            value_range = 1
        scale = int(chart_h) / value_range
        baseline_y = int(chart_top + chart_h - (-min_cumulative * scale))

        bar_total_w = int(chart_w) // n
        bar_w = int(bar_total_w * 0.65)
        bar_gap = int(bar_total_w * 0.35)

        cumulative = 0
        for i, item in enumerate(items):
            x = int(chart_left) + i * bar_total_w + bar_gap // 2
            val = item["value"]
            item_type = item.get("type", "increase" if val >= 0 else "decrease")

            if item_type == "start" or item_type == "total":
                bar_h = int(abs(val) * scale)
                bar_y = baseline_y - bar_h if val >= 0 else baseline_y
                color = self.colors.primary
                cumulative = val
            elif val >= 0:
                bar_h = int(val * scale)
                bar_y = baseline_y - int(cumulative * scale) - bar_h
                color = self.colors.positive
                cumulative += val
            else:
                bar_h = int(abs(val) * scale)
                bar_y = baseline_y - int(cumulative * scale)
                color = self.colors.negative
                cumulative += val

            bar_h = max(bar_h, Inches(0.1))

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, bar_y, bar_w, bar_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = color
            shape.line.fill.background()

            # Data label
            label_y = bar_y - Inches(0.25) if val >= 0 else bar_y + bar_h
            txb = slide.shapes.add_textbox(x, int(label_y), bar_w, Inches(0.25))
            tf = txb.text_frame
            p = tf.paragraphs[0]
            prefix = "+" if val > 0 and item_type not in ("start", "total") else ""
            p.text = f"{prefix}{val:,.0f}"
            p.font.name = self.typo.body_font
            p.font.size = Pt(9)
            p.font.color.rgb = self.colors.text_dark
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            # Category label
            cat_txb = slide.shapes.add_textbox(
                x, int(chart_top + chart_h + Inches(0.1)), bar_w, Inches(0.3)
            )
            cat_tf = cat_txb.text_frame
            cat_p = cat_tf.paragraphs[0]
            cat_p.text = item.get("label", "")
            cat_p.font.name = self.typo.body_font
            cat_p.font.size = Pt(8)
            cat_p.font.color.rgb = self.colors.text_mid
            cat_p.alignment = PP_ALIGN.CENTER

        self._style_common(slide, data)

    def _render_takeaways(self, slide, data: dict):
        """Key takeaways / conclusion slide."""
        self._add_header_bar(slide, data.get("title", "Key Takeaways"),
                             subhead=data.get("subhead"),
                             message_line=data.get("message_line") or data.get("subtitle"))
        if self.uses_template:
            self._remove_placeholder(slide, 1)
        points = data.get("points", [])
        if points:
            self._add_bullets(
                slide, points,
                MARGIN_L, CONTENT_TOP + Inches(0.3),
                CONTENT_W, Inches(3.0)
            )

        next_steps = data.get("next_steps")
        if next_steps:
            # "Next Steps" sub-header
            txb = slide.shapes.add_textbox(
                MARGIN_L, CONTENT_TOP + Inches(3.5), Inches(3.0), Inches(0.4)
            )
            tf = txb.text_frame
            p = tf.paragraphs[0]
            p.text = "Next Steps"
            p.font.name = self.typo.body_font
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.colors.primary

            self._add_bullets(
                slide, next_steps,
                MARGIN_L, CONTENT_TOP + Inches(4.0),
                CONTENT_W, Inches(2.0)
            )

        self._add_page_number(slide)

    def _render_thank_you(self, slide, data: dict):
        """Thank you / end slide."""
        if self.style == "accenture":
            if self.uses_template:
                # Section divider layout: purple bg from layout
                # Repurpose section title for "Thank you" centered
                self._remove_placeholder(slide, 10)  # no section number
                self._set_placeholder_text(slide, 0, data.get("title", "Thank you"),
                    font_name=self.typo.title_font, font_size=Pt(44),
                    bold=True, color=self.colors.white, alignment=PP_ALIGN.CENTER)
                if data.get("contact"):
                    self._set_placeholder_text(slide, 11, data["contact"],
                        font_name=self.typo.body_font, font_size=Pt(14),
                        color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)
                else:
                    self._remove_placeholder(slide, 11)
                return
            # Fallback: programmatic purple bg
            bg_rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
            )
            bg_rect.fill.solid()
            bg_rect.fill.fore_color.rgb = self.colors.primary
            bg_rect.line.fill.background()
        else:
            # McKinsey: solid primary background
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = self.colors.primary

        txb = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.5), Inches(10.0), Inches(2.0)
        )
        tf = txb.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "Thank you")
        p.font.name = self.typo.title_font
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors.white
        p.alignment = PP_ALIGN.CENTER

        if data.get("contact"):
            p2 = tf.add_paragraph()
            p2.text = data["contact"]
            p2.font.name = self.typo.body_font
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(20)

    # ── Accenture Component Helpers (050113-verified) ──

    def _add_callout_box(self, slide, text, x, y, w, h,
                         fill=True, font_size=None):
        """Add purple callout box. 050113: stroke=#A000FF, fill=#EBCCFF or transparent."""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors.accent2  # #EBCCFF
        else:
            shape.fill.background()
        shape.line.color.rgb = self.colors.primary  # #A000FF
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.typo.body_font
        p.font.size = font_size or Pt(12)
        p.font.color.rgb = self.colors.text_dark
        return shape

    def _add_info_panel(self, slide, text, x, y, w, h,
                        bg_color=None, font_size=None):
        """Add gray info panel. 050113: #D8D9D8 background box with content."""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color or self.colors.gridline  # #D8D9D8
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)
        tf.margin_top = Pt(8)
        tf.margin_bottom = Pt(8)
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.typo.body_font
        p.font.size = font_size or Pt(14)
        p.font.color.rgb = self.colors.text_dark
        return shape

    def _add_tab_bar(self, slide, tabs, active_index=None):
        """Add small colored tab bar at top-right. 050113: active=#A000FF, inactive=#BEBFBE.
        Tabs are ~48x27pt each, right-aligned from x=621."""
        n = len(tabs)
        if n == 0:
            return
        tab_w = Pt(48)
        tab_h = Pt(27)
        gap = Pt(2)
        total_w = tab_w * n + gap * (n - 1)
        start_x = Pt(925) - total_w  # right-aligned
        y = Pt(6)

        for i, label in enumerate(tabs):
            x = start_x + i * (tab_w + gap)
            is_active = (active_index is not None and i == active_index)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, tab_w, tab_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors.primary if is_active else RGBColor(0xBE, 0xBF, 0xBE)
            shape.line.fill.background()
            tf = shape.text_frame
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = label
            p.font.name = self.typo.body_font
            p.font.size = Pt(10)
            p.font.bold = is_active
            p.font.color.rgb = self.colors.white
            p.alignment = PP_ALIGN.CENTER

    def _add_step_number_icon(self, slide, number, x=None, y=None):
        """Add step number box. 050113: 28x28pt #7500C0 square with white number."""
        x = x or Pt(35)
        y = y or Pt(35)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(28), Pt(28))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors.primary_light  # #7500C0
        shape.line.fill.background()
        tf = shape.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.name = self.typo.title_font
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.colors.white
        p.alignment = PP_ALIGN.CENTER
        return shape

    # ── New Slide Renderers (050113-based) ──

    def _render_content_with_panel(self, slide, data: dict):
        """Content + Right Info Panel. 050113 pattern: main content left, gray panel right."""
        msg = data.get("message_line") or data.get("subtitle")
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=msg)
        if self.uses_template:
            self._remove_placeholder(slide, 1)
            self._remove_placeholder(slide, 13)

        body_top = Pt(155) if (self.style == "accenture" and msg) else CONTENT_TOP + Inches(0.2)
        bullets = data.get("bullets", [])
        if bullets:
            self._add_bullets(slide, bullets,
                Pt(35) if self.style == "accenture" else MARGIN_L,
                body_top, Pt(525), Pt(355))

        # Right panel
        panel = data.get("panel", {})
        panel_text = panel.get("text", "")
        if panel_text:
            self._add_info_panel(slide, panel_text,
                Pt(580), body_top, Pt(345), Pt(355))

        # Callout on panel (optional)
        callout = data.get("callout")
        if callout:
            self._add_callout_box(slide, callout,
                Pt(580), Pt(400), Pt(345), Pt(100))

        self._style_common(slide, data)

    def _render_tab_content(self, slide, data: dict):
        """Content with tab bar navigation. 050113 pattern: small tabs at top-right."""
        msg = data.get("message_line") or data.get("subtitle")
        self._add_header_bar(slide, data.get("title", ""),
                             subhead=data.get("subhead"),
                             message_line=msg)
        if self.uses_template:
            self._remove_placeholder(slide, 1)
            self._remove_placeholder(slide, 18)

        # Draw tab bar
        tabs = data.get("tabs", [])
        active_tab = data.get("active_tab")
        if tabs:
            self._add_tab_bar(slide, tabs, active_tab)

        body_top = Pt(155) if (self.style == "accenture" and msg) else CONTENT_TOP + Inches(0.2)
        bullets = data.get("bullets", [])
        if bullets:
            body_h = Pt(355) if self.style == "accenture" else CONTENT_H - Inches(0.5)
            self._add_bullets(slide, bullets,
                Pt(35) if self.style == "accenture" else MARGIN_L,
                body_top, CONTENT_W, body_h)

        self._style_common(slide, data)

    def _render_step_content(self, slide, data: dict):
        """Numbered step content. 050113 pattern: step box top-left, offset title."""
        msg = data.get("message_line") or data.get("subtitle")

        if self.uses_template:
            self._remove_placeholder(slide, 19)
            self._remove_placeholder(slide, 1)

        # Step number icon
        step_num = data.get("step_number", 1)
        self._add_step_number_icon(slide, step_num)

        if self.uses_template:
            # Title at offset position (right of step box)
            self._set_placeholder_text(slide, 0, data.get("title", ""),
                font_name=self.typo.title_font, font_size=Pt(28),
                bold=True, color=self.colors.text_dark)
            if msg:
                self._set_placeholder_text(slide, 11, msg,
                    font_name=self.typo.title_font, font_size=self.typo.message_size,
                    bold=True, color=self.colors.text_dark)
            else:
                self._remove_placeholder(slide, 11)
        else:
            # Programmatic title (offset right of step box)
            txb = slide.shapes.add_textbox(Pt(70), Pt(32), Inches(11), Pt(40))
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data.get("title", "")
            p.font.name = self.typo.title_font
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.colors.text_dark
            if msg:
                msg_txb = slide.shapes.add_textbox(Pt(35), Pt(80), Inches(12), Pt(44))
                msg_tf = msg_txb.text_frame
                msg_p = msg_tf.paragraphs[0]
                msg_p.text = msg
                msg_p.font.name = self.typo.title_font
                msg_p.font.size = self.typo.message_size
                msg_p.font.bold = True
                msg_p.font.color.rgb = self.colors.text_dark

        body_top = Pt(155)
        bullets = data.get("bullets", [])
        if bullets:
            self._add_bullets(slide, bullets, Pt(35), body_top, CONTENT_W, Pt(355))

        self._style_common(slide, data)

    def _render_appendix_divider(self, slide, data: dict):
        """Appendix section divider."""
        data["section_title"] = data.get("title", "Appendix")
        data.setdefault("section_number", None)
        self._render_section_divider(slide, data)


# ============================================================
# CLI Entry Point
# ============================================================

def main():
    if len(sys.argv) < 2:
        print("Usage: generate_pptx.py <input.json> [output.pptx] [--style mckinsey|accenture]")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 and not sys.argv[2].startswith("--") else None

    style = "mckinsey"
    for i, arg in enumerate(sys.argv):
        if arg == "--style" and i + 1 < len(sys.argv):
            style = sys.argv[i + 1]

    with open(input_path, "r", encoding="utf-8") as f:
        deck = json.load(f)

    if output_path is None:
        output_path = Path(input_path).stem + ".pptx"

    # CLI --style flag takes precedence; fall back to JSON "style" field
    cli_style_set = any(arg == "--style" for arg in sys.argv)
    effective_style = style if cli_style_set else deck.get("style", style)
    builder = PptxBuilder(style=effective_style)
    result = builder.build(deck, output_path)
    print(f"Generated: {result}")


if __name__ == "__main__":
    main()
