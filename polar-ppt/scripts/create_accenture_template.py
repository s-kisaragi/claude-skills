#!/usr/bin/env python3
"""
Create Accenture slide master template (.pptx) based on 050113.pdf analysis.
Layouts: Cover, Content, TwoColumn, Agenda, SectionDivider, BlankWithHeader

Zones (050113 verified):
  Subhead    Y=12pt   MeiryoUI-Bold 18pt
  Title      Y=32pt   MeiryoUI-Bold 32pt
  Message    Y=73pt   MeiryoUI-Bold 18pt
  Separator  Y=138pt  #000000 0.5pt
  Body       Y=155pt  (variable)
  Footer     Y=519pt  Graphik 8pt italic #919191
"""

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Pt, Inches, Emu
from pathlib import Path

# === Constants (050113 verified) ===
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Meiryo UI"

# Zone Y positions (in EMU)
Y_SUBHEAD = Pt(12)
Y_TITLE = Pt(32)
Y_MESSAGE = Pt(73)
Y_SEPARATOR = Pt(138)
Y_BODY = Pt(155)
Y_FOOTER = Pt(519)

X_LEFT = Pt(35)
CONTENT_W = Pt(890)  # 960 - 35*2
BODY_H = Pt(355)  # 510 - 155

# IDs: must be unique within each spTree; use high numbers to avoid conflicts
_id_counter = 100


def _next_id():
    global _id_counter
    _id_counter += 1
    return _id_counter


def _make_ph(name, ph_type, idx, x, y, cx, cy, font=FONT, sz=1800, bold=False, color="000000", anchor="t"):
    """Create a placeholder shape XML string."""
    type_attr = f'type="{ph_type}" ' if ph_type else ""
    b_attr = 'b="1"' if bold else ""
    return f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="{_next_id()}" name="{name}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph {type_attr}idx="{idx}"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{x}" y="{y}"/>
      <a:ext cx="{cx}" cy="{cy}"/>
    </a:xfrm>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" anchor="{anchor}">
      <a:noAutofit/>
    </a:bodyPr>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="ja-JP" sz="{sz}" {b_attr} dirty="0">
      <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
      <a:latin typeface="{font}"/><a:ea typeface="{font}"/>
    </a:endParaRPr></a:p>
  </p:txBody>
</p:sp>'''


def _make_textbox(name, x, y, cx, cy, text, font=FONT, sz=800, color="919191", bold=False, italic=False):
    """Create a non-placeholder textbox shape XML string."""
    b_attr = 'b="1"' if bold else ""
    i_attr = 'i="1"' if italic else ""
    return f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{_next_id()}" name="{name}"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{x}" y="{y}"/>
      <a:ext cx="{cx}" cy="{cy}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0"/>
    <a:lstStyle/>
    <a:p>
      <a:r>
        <a:rPr lang="en-US" sz="{sz}" {b_attr} {i_attr} dirty="0">
          <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
          <a:latin typeface="{font}"/><a:ea typeface="{font}"/>
        </a:rPr>
        <a:t>{text}</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>'''


def _make_line(name, x, y, cx, cy, color="000000", w=6350):
    """Create a thin line (rectangle with fill, no border)."""
    return f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{_next_id()}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{x}" y="{y}"/>
      <a:ext cx="{cx}" cy="{cy}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr>
</p:sp>'''


def _clear_layout(layout):
    """Remove all existing placeholders from a layout."""
    cSld = layout._element.find(qn('p:cSld'))
    spTree = cSld.find(qn('p:spTree'))
    for sp in list(spTree.findall(qn('p:sp'))):
        spTree.remove(sp)
    return spTree


def _add_to(spTree, xml_str):
    """Parse XML string and append to shape tree."""
    spTree.append(parse_xml(xml_str))


def setup_master_footer(slide_master):
    """Clean default shapes and add shared footer to slide master."""
    cSld = slide_master._element.find(qn('p:cSld'))
    spTree = cSld.find(qn('p:spTree'))

    # Remove all default master shapes (Title, Text, Date, Footer, SlideNum placeholders)
    for sp in list(spTree.findall(qn('p:sp'))):
        spTree.remove(sp)

    # Copyright text — bottom-left, Graphik-style 8pt italic #919191
    _add_to(spTree, _make_textbox(
        "Master Copyright", X_LEFT, Y_FOOTER, Inches(6), Pt(12),
        "Copyright \u00a9 2026 Accenture. All rights reserved.",
        font=FONT, sz=800, color="919191", italic=True
    ))


def setup_cover_layout(layout):
    """Layout 0: Cover/Title slide."""
    layout._element.find(qn('p:cSld')).set('name', 'Accenture Cover')
    spTree = _clear_layout(layout)

    # Addressee — idx=14
    _add_to(spTree, _make_ph("Addressee", "body", 14,
        X_LEFT, Pt(119), Inches(8), Pt(22),
        sz=1600, bold=False))

    # Project subtitle — idx=15
    _add_to(spTree, _make_ph("Project Subtitle", "body", 15,
        X_LEFT, Pt(157), Inches(11), Pt(24),
        sz=1800, bold=True))

    # Main title — idx=0
    _add_to(spTree, _make_ph("Title", "ctrTitle", 0,
        X_LEFT, Pt(190), Inches(11), Pt(70),
        sz=3200, bold=True))

    # Document type — idx=16
    _add_to(spTree, _make_ph("Document Type", "body", 16,
        X_LEFT, Pt(284), Inches(5), Pt(28),
        sz=2000, bold=True))

    # Date — idx=17
    _add_to(spTree, _make_ph("Date", "body", 17,
        Pt(80), Pt(373), Inches(4), Pt(28),
        sz=2000, bold=True))


def setup_content_layout(layout):
    """Layout 1: Standard content slide (subhead/title/message/body/footer)."""
    layout._element.find(qn('p:cSld')).set('name', 'Accenture Content')
    spTree = _clear_layout(layout)

    # Subhead — idx=10
    _add_to(spTree, _make_ph("Subhead", "body", 10,
        X_LEFT, Y_SUBHEAD, Inches(12), Pt(22),
        sz=1800, bold=True))

    # Title — idx=0
    _add_to(spTree, _make_ph("Title", "title", 0,
        X_LEFT, Y_TITLE, Inches(12), Pt(40),
        sz=3200, bold=True))

    # Message line — idx=11
    _add_to(spTree, _make_ph("Message", "body", 11,
        X_LEFT, Y_MESSAGE, Inches(12), Pt(44),
        sz=1800, bold=True))

    # Separator line (non-placeholder, on layout)
    _add_to(spTree, _make_line("Separator", X_LEFT, Y_SEPARATOR, Pt(890), Pt(1)))

    # Body — idx=1
    _add_to(spTree, _make_ph("Body", "body", 1,
        X_LEFT, Y_BODY, CONTENT_W, BODY_H,
        sz=1600, bold=False))

    # Page number placeholder — idx=12
    _add_to(spTree, _make_ph("Page Number", "sldNum", 12,
        Pt(890), Y_FOOTER, Pt(40), Pt(12),
        sz=800, color="919191"))


def setup_twocol_layout(layout):
    """Layout 3: Two-column content slide."""
    layout._element.find(qn('p:cSld')).set('name', 'Accenture Two Column')
    spTree = _clear_layout(layout)

    # Same header zones as content
    _add_to(spTree, _make_ph("Subhead", "body", 10,
        X_LEFT, Y_SUBHEAD, Inches(12), Pt(22),
        sz=1800, bold=True))

    _add_to(spTree, _make_ph("Title", "title", 0,
        X_LEFT, Y_TITLE, Inches(12), Pt(40),
        sz=3200, bold=True))

    _add_to(spTree, _make_ph("Message", "body", 11,
        X_LEFT, Y_MESSAGE, Inches(12), Pt(44),
        sz=1800, bold=True))

    _add_to(spTree, _make_line("Separator", X_LEFT, Y_SEPARATOR, Pt(890), Pt(1)))

    # Left body — idx=1
    half_w = int(Pt(430))
    _add_to(spTree, _make_ph("Left Body", "body", 1,
        X_LEFT, Y_BODY, half_w, BODY_H,
        sz=1600, bold=False))

    # Right body — idx=13
    _add_to(spTree, _make_ph("Right Body", "body", 13,
        Pt(35 + 445), Y_BODY, half_w, BODY_H,
        sz=1600, bold=False))

    # Page number
    _add_to(spTree, _make_ph("Page Number", "sldNum", 12,
        Pt(890), Y_FOOTER, Pt(40), Pt(12),
        sz=800, color="919191"))


def setup_agenda_layout(layout):
    """Layout 5: Agenda slide."""
    layout._element.find(qn('p:cSld')).set('name', 'Accenture Agenda')
    spTree = _clear_layout(layout)

    # Title — idx=0 (agenda title e.g. "アジェンダ")
    _add_to(spTree, _make_ph("Title", "title", 0,
        X_LEFT, Pt(32), Inches(12), Pt(36),
        sz=2800, bold=True))

    # Body area for agenda items — idx=1
    _add_to(spTree, _make_ph("Agenda Items", "body", 1,
        Pt(0), Pt(160), Pt(960), Pt(340),
        sz=2400, bold=False))

    # Page number
    _add_to(spTree, _make_ph("Page Number", "sldNum", 12,
        Pt(890), Y_FOOTER, Pt(40), Pt(12),
        sz=800, color="919191"))


def setup_divider_layout(layout):
    """Layout 2: Section divider (purple background)."""
    layout._element.find(qn('p:cSld')).set('name', 'Accenture Section Divider')
    spTree = _clear_layout(layout)

    # Full-slide purple background
    _add_to(spTree, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{_next_id()}" name="Purple BG"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="0" y="0"/>
      <a:ext cx="{SLIDE_W}" cy="{SLIDE_H}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="A000FF"/></a:solidFill>
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr>
</p:sp>''')

    # Section number — idx=10
    _add_to(spTree, _make_ph("Section Number", "body", 10,
        Pt(48), Pt(180), Pt(80), Pt(50),
        sz=4000, bold=True, color="FFFFFF"))

    # Section title — idx=0
    _add_to(spTree, _make_ph("Section Title", "title", 0,
        Pt(48), Pt(240), Pt(860), Pt(80),
        sz=3200, bold=True, color="FFFFFF"))

    # Subtitle — idx=11
    _add_to(spTree, _make_ph("Subtitle", "body", 11,
        Pt(48), Pt(330), Pt(860), Pt(40),
        sz=1800, bold=False, color="DDDDDD"))


def setup_content_right_panel_layout(layout):
    """Layout 4: Content + Right Info Panel (60/40 split, gray panel on right).
    050113 pattern: pages 11,25,28,35,103,140 — main content left, #d8d9d8 panel right."""
    layout._element.find(qn('p:cSld')).set('name', 'Accenture Content + Panel')
    spTree = _clear_layout(layout)

    # Header zones (same as content)
    _add_to(spTree, _make_ph("Subhead", "body", 10,
        X_LEFT, Y_SUBHEAD, Inches(12), Pt(22),
        sz=1800, bold=True))
    _add_to(spTree, _make_ph("Title", "title", 0,
        X_LEFT, Y_TITLE, Inches(12), Pt(40),
        sz=3200, bold=True))
    _add_to(spTree, _make_ph("Message", "body", 11,
        X_LEFT, Y_MESSAGE, Inches(12), Pt(44),
        sz=1800, bold=True))
    _add_to(spTree, _make_line("Separator", X_LEFT, Y_SEPARATOR, Pt(890), Pt(1)))

    # Left body — idx=1 (60% width, ~530pt)
    _add_to(spTree, _make_ph("Left Body", "body", 1,
        X_LEFT, Y_BODY, Pt(530), BODY_H,
        sz=1600, bold=False))

    # Right panel background — gray #D8D9D8 (non-placeholder visual)
    _add_to(spTree, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{_next_id()}" name="Panel BG"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{Pt(580)}" y="{Y_BODY}"/>
      <a:ext cx="{Pt(345)}" cy="{BODY_H}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="D8D9D8"/></a:solidFill>
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr>
</p:sp>''')

    # Right panel content — idx=13
    _add_to(spTree, _make_ph("Panel Content", "body", 13,
        Pt(590), Pt(165), Pt(325), Pt(340),
        sz=1400, bold=False))

    # Page number
    _add_to(spTree, _make_ph("Page Number", "sldNum", 12,
        Pt(890), Y_FOOTER, Pt(40), Pt(12),
        sz=800, color="919191"))


def setup_tab_bar_layout(layout):
    """Layout 7: Content with Tab Bar (small colored tabs at top-right).
    050113 pattern: pages 10-36 — 6 small tabs at top-right, active=#a000ff, inactive=#bebfbe."""
    layout._element.find(qn('p:cSld')).set('name', 'Accenture Tab Bar Content')
    spTree = _clear_layout(layout)

    # Tab bar zone — idx=18 (placeholder for tab labels, rendered programmatically)
    # Tabs sit at y=6-34, right-aligned (x=621-921)
    _add_to(spTree, _make_ph("Tab Bar", "body", 18,
        Pt(620), Pt(6), Pt(305), Pt(28),
        sz=1200, bold=False, color="919191"))

    # Header zones (subhead/title/message) — same positions
    _add_to(spTree, _make_ph("Subhead", "body", 10,
        X_LEFT, Y_SUBHEAD, Inches(8), Pt(22),
        sz=1800, bold=True))
    _add_to(spTree, _make_ph("Title", "title", 0,
        X_LEFT, Y_TITLE, Inches(12), Pt(40),
        sz=3200, bold=True))
    _add_to(spTree, _make_ph("Message", "body", 11,
        X_LEFT, Y_MESSAGE, Inches(12), Pt(44),
        sz=1800, bold=True))
    _add_to(spTree, _make_line("Separator", X_LEFT, Y_SEPARATOR, Pt(890), Pt(1)))

    # Body — idx=1
    _add_to(spTree, _make_ph("Body", "body", 1,
        X_LEFT, Y_BODY, CONTENT_W, BODY_H,
        sz=1600, bold=False))

    # Page number
    _add_to(spTree, _make_ph("Page Number", "sldNum", 12,
        Pt(890), Y_FOOTER, Pt(40), Pt(12),
        sz=800, color="919191"))


def setup_step_content_layout(layout):
    """Layout 8: Numbered Step Content (step icon top-left, offset title).
    050113 pattern: pages 138-145 — 28x28pt #7500C0 step box, title offset to right."""
    layout._element.find(qn('p:cSld')).set('name', 'Accenture Step Content')
    spTree = _clear_layout(layout)

    # Step number box background — #7500C0, 28x28pt at (35, 35)
    _add_to(spTree, f'''<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="{_next_id()}" name="Step Box BG"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{X_LEFT}" y="{X_LEFT}"/>
      <a:ext cx="{Pt(28)}" cy="{Pt(28)}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="7500C0"/></a:solidFill>
    <a:ln w="0"><a:noFill/></a:ln>
  </p:spPr>
</p:sp>''')

    # Step number text — idx=19
    _add_to(spTree, _make_ph("Step Number", "body", 19,
        X_LEFT, X_LEFT, Pt(28), Pt(28),
        sz=1800, bold=True, color="FFFFFF", anchor="ctr"))

    # Title — idx=0, offset right of step box (x=70)
    _add_to(spTree, _make_ph("Title", "title", 0,
        Pt(70), Pt(32), Inches(11), Pt(40),
        sz=2800, bold=True))

    # Message — idx=11
    _add_to(spTree, _make_ph("Message", "body", 11,
        X_LEFT, Pt(80), Inches(12), Pt(44),
        sz=1800, bold=True))

    _add_to(spTree, _make_line("Separator", X_LEFT, Pt(138), Pt(890), Pt(1)))

    # Body — idx=1
    _add_to(spTree, _make_ph("Body", "body", 1,
        X_LEFT, Y_BODY, CONTENT_W, BODY_H,
        sz=1600, bold=False))

    # Page number
    _add_to(spTree, _make_ph("Page Number", "sldNum", 12,
        Pt(890), Y_FOOTER, Pt(40), Pt(12),
        sz=800, color="919191"))


def setup_blank_layout(layout):
    """Layout 6: Blank with header zones (for charts, custom content)."""
    layout._element.find(qn('p:cSld')).set('name', 'Accenture Blank With Header')
    spTree = _clear_layout(layout)

    # Same header zones
    _add_to(spTree, _make_ph("Subhead", "body", 10,
        X_LEFT, Y_SUBHEAD, Inches(12), Pt(22),
        sz=1800, bold=True))

    _add_to(spTree, _make_ph("Title", "title", 0,
        X_LEFT, Y_TITLE, Inches(12), Pt(40),
        sz=3200, bold=True))

    _add_to(spTree, _make_ph("Message", "body", 11,
        X_LEFT, Y_MESSAGE, Inches(12), Pt(44),
        sz=1800, bold=True))

    _add_to(spTree, _make_line("Separator", X_LEFT, Y_SEPARATOR, Pt(890), Pt(1)))

    # Page number
    _add_to(spTree, _make_ph("Page Number", "sldNum", 12,
        Pt(890), Y_FOOTER, Pt(40), Pt(12),
        sz=800, color="919191"))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Setup slide master with shared footer
    setup_master_footer(prs.slide_masters[0])

    # Repurpose default layouts:
    # 0 = Cover, 1 = Content, 2 = SectionDivider, 3 = TwoColumn,
    # 4 = Content+Panel, 5 = Agenda, 6 = BlankWithHeader,
    # 7 = TabBarContent, 8 = StepContent
    setup_cover_layout(prs.slide_layouts[0])                # Cover
    setup_content_layout(prs.slide_layouts[1])               # Content
    setup_divider_layout(prs.slide_layouts[2])               # Section Divider
    setup_twocol_layout(prs.slide_layouts[3])                # Two Column
    setup_content_right_panel_layout(prs.slide_layouts[4])   # Content + Panel
    setup_agenda_layout(prs.slide_layouts[5])                # Agenda
    setup_blank_layout(prs.slide_layouts[6])                 # Blank + Header
    setup_tab_bar_layout(prs.slide_layouts[7])               # Tab Bar Content
    setup_step_content_layout(prs.slide_layouts[8])          # Step Content

    # Save template
    output_dir = Path(__file__).parent.parent / "assets"
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / "accenture_template.pptx"
    prs.save(str(output_path))
    print(f"Template saved: {output_path}")


if __name__ == "__main__":
    main()
